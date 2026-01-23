"""
Document Processor for handling various file formats
"""
import PyPDF2
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from PIL import Image
import pytesseract
from typing import Dict, List
import os
import io


class DocumentProcessor:
    """Process and extract text from various document formats"""
    
    @staticmethod
    def process_file(file_path: str, file_content=None) -> Dict[str, any]:
        """Process uploaded file and extract text content"""
        ext = os.path.splitext(file_path)[-1].lower()
        
        result = {
            "filename": os.path.basename(file_path),
            "extension": ext,
            "text": "",
            "metadata": {},
            "chunks": []
        }
        
        try:
            if ext == ".pdf":
                result["text"] = DocumentProcessor.extract_from_pdf(file_content)
            elif ext in [".docx", ".doc"]:
                result["text"] = DocumentProcessor.extract_from_docx(file_content)
            elif ext in [".pptx", ".ppt"]:
                result["text"] = DocumentProcessor.extract_from_pptx(file_content)
            elif ext in [".xlsx", ".xls"]:
                result["text"] = DocumentProcessor.extract_from_xlsx(file_content)
            elif ext in [".png", ".jpg", ".jpeg", ".tiff", ".bmp"]:
                result["text"] = DocumentProcessor.extract_from_image(file_content)
            elif ext in [".txt", ".md"]:
                result["text"] = file_content.decode("utf-8") if isinstance(file_content, bytes) else file_content.read().decode("utf-8")
            else:
                result["text"] = ""
            
            # Chunk the text for better processing
            result["chunks"] = DocumentProcessor.chunk_text(result["text"])
            result["metadata"]["word_count"] = len(result["text"].split())
            result["metadata"]["chunk_count"] = len(result["chunks"])
            result["metadata"]["source_file"] = file_path
            
        except Exception as e:
            result["error"] = str(e)
        
        return result
    
    @staticmethod
    def extract_from_pdf(file_content) -> str:
        """Extract text from PDF file"""
        text = ""
        try:
            pdf_reader = PyPDF2.PdfReader(file_content)
            for i, page in enumerate(pdf_reader.pages):
                page_text = page.extract_text()
                text += f"\n[Page {i+1}]\n{page_text}\n"
        except Exception as e:
            text = f"Error extracting PDF: {str(e)}"
        return text
    
    @staticmethod
    def extract_from_docx(file_content) -> str:
        """Extract text from DOCX file"""
        text = ""
        try:
            doc = Document(file_content)
            for i, paragraph in enumerate(doc.paragraphs):
                text += paragraph.text + "\n"
        except Exception as e:
            text = f"Error extracting DOCX: {str(e)}"
        return text
    
    @staticmethod
    def extract_from_pptx(file_content) -> str:
        """Extract text from PPTX file"""
        text = ""
        try:
            prs = Presentation(file_content)
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"\n[Slide {slide_num}]\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            text = f"Error extracting PPTX: {str(e)}"
        return text
    
    @staticmethod
    def extract_from_xlsx(file_content) -> str:
        """Extract text from XLSX file"""
        text = ""
        try:
            # Load workbook
            wb = load_workbook(file_content, read_only=True)
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text += f"\n[Sheet: {sheet_name}]\n"
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        text += row_text + "\n"
        except Exception as e:
            text = f"Error extracting XLSX: {str(e)}"
        return text
    
    @staticmethod
    def extract_from_image(file_content) -> str:
        """Extract text from image using OCR"""
        text = ""
        try:
            # Open image
            image = Image.open(file_content)
            # Perform OCR
            text = pytesseract.image_to_string(image)
            if not text.strip():
                text = "[No text detected in image]"
        except Exception as e:
            text = f"Error extracting from image (OCR): {str(e)}\nNote: Ensure Tesseract OCR is installed on your system."
        return text
    
    @staticmethod
    def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """Split text into overlapping chunks"""
        words = text.split()
        chunks = []
        
        for i in range(0, len(words), chunk_size - overlap):
            chunk = " ".join(words[i:i + chunk_size])
            if chunk:
                chunks.append(chunk)
        
        return chunks if chunks else [text]
